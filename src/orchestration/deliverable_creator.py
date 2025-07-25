"""
Deliverable Creator Module for AgentFlow
Creates actual files and deliverables based on AI-generated content
"""
import os
import json
import logging
from datetime import datetime
from typing import Dict, List, Any, Optional
from pathlib import Path

# Import libraries for file creation
try:
    from fpdf import FPDF
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    # Install required packages if not available
    import subprocess
    import sys
    
    packages = ['fpdf2', 'python-pptx']
    for package in packages:
        try:
            subprocess.check_call([sys.executable, '-m', 'pip', 'install', package])
        except:
            pass
    
    # Try importing again
    try:
        from fpdf import FPDF
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
    except ImportError:
        # Fallback to basic file creation
        FPDF = None
        Presentation = None

logger = logging.getLogger(__name__)

class DeliverableCreator:
    """Creates actual deliverable files from AI-generated content"""
    
    def __init__(self, base_output_dir: str = "/tmp/agentflow_deliverables"):
        self.base_output_dir = Path(base_output_dir)
        self.base_output_dir.mkdir(parents=True, exist_ok=True)
        
        # Create subdirectories
        (self.base_output_dir / "presentations").mkdir(exist_ok=True)
        (self.base_output_dir / "documents").mkdir(exist_ok=True)
        (self.base_output_dir / "reports").mkdir(exist_ok=True)
        (self.base_output_dir / "data").mkdir(exist_ok=True)
        
        logger.info(f"📁 Deliverable Creator initialized with output directory: {self.base_output_dir}")
    
    def create_powerpoint_presentation(self, content: str, task: str, execution_id: str) -> str:
        """Create an actual PowerPoint presentation file"""
        try:
            if not Presentation:
                # Fallback to HTML presentation
                return self._create_html_presentation(content, task, execution_id)
            
            logger.info("🎯 Creating PowerPoint presentation...")
            
            # Create presentation
            prs = Presentation()
            
            # Parse content to extract slides
            slides_content = self._parse_content_for_slides(content, task)
            
            # Create title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            # Extract title from task or use default
            presentation_title = self._extract_title_from_task(task)
            title.text = presentation_title
            subtitle.text = f"Generated by AgentFlow AI\n{datetime.now().strftime('%B %d, %Y')}"
            
            # Create content slides
            for i, slide_content in enumerate(slides_content):
                self._add_content_slide(prs, slide_content, i + 1)
            
            # Save presentation
            filename = f"AI_Presentation_{execution_id}.pptx"
            filepath = self.base_output_dir / "presentations" / filename
            prs.save(str(filepath))
            
            logger.info(f"✅ PowerPoint presentation created: {filepath}")
            return str(filepath)
            
        except Exception as e:
            logger.error(f"❌ Failed to create PowerPoint: {e}")
            # Fallback to HTML presentation
            return self._create_html_presentation(content, task, execution_id)
    
    def create_pdf_document(self, content: str, task: str, execution_id: str) -> str:
        """Create an actual PDF document"""
        try:
            if not FPDF:
                # Fallback to text document
                return self._create_text_document(content, task, execution_id)
            
            logger.info("📄 Creating PDF document...")
            
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font('Arial', 'B', 16)
            
            # Title
            title = self._extract_title_from_task(task)
            pdf.cell(0, 10, title, ln=True, align='C')
            pdf.ln(10)
            
            # Content
            pdf.set_font('Arial', '', 12)
            
            # Split content into paragraphs and add to PDF
            paragraphs = content.split('\n\n')
            for paragraph in paragraphs:
                if paragraph.strip():
                    # Handle long paragraphs
                    words = paragraph.split(' ')
                    line = ""
                    for word in words:
                        if len(line + word) < 80:  # Approximate character limit per line
                            line += word + " "
                        else:
                            pdf.cell(0, 6, line.strip(), ln=True)
                            line = word + " "
                    if line.strip():
                        pdf.cell(0, 6, line.strip(), ln=True)
                    pdf.ln(3)
            
            # Save PDF
            filename = f"AI_Document_{execution_id}.pdf"
            filepath = self.base_output_dir / "documents" / filename
            pdf.output(str(filepath))
            
            logger.info(f"✅ PDF document created: {filepath}")
            return str(filepath)
            
        except Exception as e:
            logger.error(f"❌ Failed to create PDF: {e}")
            # Fallback to text document
            return self._create_text_document(content, task, execution_id)
    
    def create_research_report(self, content: str, task: str, execution_id: str) -> str:
        """Create a structured research report"""
        try:
            logger.info("📊 Creating research report...")
            
            # Create structured report content
            report_content = self._structure_research_report(content, task)
            
            # Try to create PDF first, fallback to text
            try:
                return self._create_pdf_report(report_content, task, execution_id)
            except:
                return self._create_text_report(report_content, task, execution_id)
                
        except Exception as e:
            logger.error(f"❌ Failed to create research report: {e}")
            return self._create_text_document(content, task, execution_id)
    
    def create_data_analysis_report(self, content: str, task: str, execution_id: str) -> str:
        """Create a data analysis report with visualizations"""
        try:
            logger.info("📈 Creating data analysis report...")
            
            # Create analysis report
            report_content = self._structure_data_analysis(content, task)
            
            # Create HTML report with potential for charts
            return self._create_html_analysis_report(report_content, task, execution_id)
            
        except Exception as e:
            logger.error(f"❌ Failed to create data analysis report: {e}")
            return self._create_text_document(content, task, execution_id)
    
    def create_deliverables_from_task(self, content: str, task: str, execution_id: str) -> List[str]:
        """Create appropriate deliverables based on task type"""
        deliverables = []
        task_lower = task.lower()
        
        try:
            # Determine deliverable types based on task
            if any(keyword in task_lower for keyword in ['powerpoint', 'presentation', 'slides', 'ppt']):
                ppt_file = self.create_powerpoint_presentation(content, task, execution_id)
                deliverables.append(ppt_file)
                
                # Also create speaker notes as separate document
                notes_file = self._create_speaker_notes(content, task, execution_id)
                deliverables.append(notes_file)
            
            elif any(keyword in task_lower for keyword in ['research', 'analysis', 'study', 'report']):
                if 'data' in task_lower or 'analysis' in task_lower:
                    report_file = self.create_data_analysis_report(content, task, execution_id)
                else:
                    report_file = self.create_research_report(content, task, execution_id)
                deliverables.append(report_file)
            
            elif any(keyword in task_lower for keyword in ['document', 'write', 'create', 'generate']):
                doc_file = self.create_pdf_document(content, task, execution_id)
                deliverables.append(doc_file)
            
            else:
                # Default: create both a document and summary
                doc_file = self.create_pdf_document(content, task, execution_id)
                deliverables.append(doc_file)
            
            # Always create a JSON metadata file
            metadata_file = self._create_metadata_file(content, task, execution_id, deliverables)
            deliverables.append(metadata_file)
            
            logger.info(f"✅ Created {len(deliverables)} deliverable files")
            return deliverables
            
        except Exception as e:
            logger.error(f"❌ Failed to create deliverables: {e}")
            # Fallback: create at least a text file
            fallback_file = self._create_text_document(content, task, execution_id)
            return [fallback_file]
    
    # Helper methods for content parsing and structuring
    
    def _parse_content_for_slides(self, content: str, task: str) -> List[Dict[str, str]]:
        """Parse content into slide structure"""
        slides = []
        
        # Try to identify slide structure in content
        if "slide" in content.lower():
            # Content already has slide structure
            slide_sections = content.split("slide")
            for i, section in enumerate(slide_sections[1:], 1):  # Skip first empty split
                title = f"Slide {i}"
                content_text = section.strip()
                
                # Extract title if present
                lines = content_text.split('\n')
                if lines and len(lines[0]) < 100:
                    title = lines[0].strip(':').strip()
                    content_text = '\n'.join(lines[1:])
                
                slides.append({
                    "title": title,
                    "content": content_text[:500]  # Limit content per slide
                })
        else:
            # Create slides from content sections
            paragraphs = content.split('\n\n')
            current_slide = {"title": "Introduction", "content": ""}
            
            for paragraph in paragraphs:
                if len(current_slide["content"]) > 300:  # Start new slide
                    slides.append(current_slide)
                    current_slide = {"title": f"Slide {len(slides) + 2}", "content": paragraph}
                else:
                    current_slide["content"] += "\n\n" + paragraph
            
            if current_slide["content"]:
                slides.append(current_slide)
        
        # Ensure we have at least 3 slides
        while len(slides) < 3:
            slides.append({
                "title": f"Additional Content {len(slides) + 1}",
                "content": "Additional information and details will be provided here."
            })
        
        return slides[:10]  # Limit to 10 slides max
    
    def _add_content_slide(self, prs, slide_content: Dict[str, str], slide_num: int):
        """Add a content slide to presentation"""
        try:
            # Use title and content layout
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            title = slide.shapes.title
            title.text = slide_content["title"]
            
            # Set content
            content_placeholder = slide.placeholders[1]
            content_placeholder.text = slide_content["content"]
            
            # Format text
            for paragraph in content_placeholder.text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.name = 'Arial'
                
        except Exception as e:
            logger.warning(f"⚠️ Failed to format slide {slide_num}: {e}")
    
    def _extract_title_from_task(self, task: str) -> str:
        """Extract a title from the task description"""
        # Look for quoted titles
        import re
        quoted_match = re.search(r'"([^"]+)"', task)
        if quoted_match:
            return quoted_match.group(1)
        
        # Look for "about X" patterns
        about_match = re.search(r'about\s+([^.]+)', task, re.IGNORECASE)
        if about_match:
            return about_match.group(1).strip()
        
        # Default title
        if 'powerpoint' in task.lower() or 'presentation' in task.lower():
            return "AI-Generated Presentation"
        elif 'research' in task.lower():
            return "Research Report"
        elif 'analysis' in task.lower():
            return "Analysis Report"
        else:
            return "AI-Generated Document"
    
    def _create_html_presentation(self, content: str, task: str, execution_id: str) -> str:
        """Create HTML presentation as fallback"""
        try:
            logger.info("🌐 Creating HTML presentation...")
            
            title = self._extract_title_from_task(task)
            slides_content = self._parse_content_for_slides(content, task)
            
            html_content = f"""
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <style>
        body {{ font-family: Arial, sans-serif; margin: 0; padding: 20px; background: #f5f5f5; }}
        .slide {{ background: white; margin: 20px 0; padding: 30px; border-radius: 8px; box-shadow: 0 2px 10px rgba(0,0,0,0.1); }}
        .slide h1 {{ color: #2c3e50; border-bottom: 3px solid #3498db; padding-bottom: 10px; }}
        .slide h2 {{ color: #34495e; }}
        .slide p {{ line-height: 1.6; color: #555; }}
        .title-slide {{ text-align: center; background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; }}
        .title-slide h1 {{ border-bottom: none; color: white; font-size: 2.5em; }}
        .slide-number {{ position: absolute; top: 10px; right: 20px; color: #888; }}
    </style>
</head>
<body>
    <div class="slide title-slide">
        <h1>{title}</h1>
        <p>Generated by AgentFlow AI</p>
        <p>{datetime.now().strftime('%B %d, %Y')}</p>
    </div>
"""
            
            for i, slide in enumerate(slides_content):
                html_content += f"""
    <div class="slide">
        <div class="slide-number">Slide {i + 2}</div>
        <h2>{slide['title']}</h2>
        <p>{slide['content'].replace(chr(10), '<br>')}</p>
    </div>
"""
            
            html_content += """
</body>
</html>
"""
            
            filename = f"AI_Presentation_{execution_id}.html"
            filepath = self.base_output_dir / "presentations" / filename
            
            with open(filepath, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            logger.info(f"✅ HTML presentation created: {filepath}")
            return str(filepath)
            
        except Exception as e:
            logger.error(f"❌ Failed to create HTML presentation: {e}")
            return self._create_text_document(content, task, execution_id)
    
    def _create_text_document(self, content: str, task: str, execution_id: str) -> str:
        """Create a simple text document as ultimate fallback"""
        try:
            logger.info("📝 Creating text document...")
            
            title = self._extract_title_from_task(task)
            
            doc_content = f"""
{title}
{'=' * len(title)}

Generated by AgentFlow AI
Date: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}
Task: {task}

{'-' * 50}

{content}

{'-' * 50}

This document was automatically generated by AgentFlow's AI system.
Execution ID: {execution_id}
"""
            
            filename = f"AI_Document_{execution_id}.txt"
            filepath = self.base_output_dir / "documents" / filename
            
            with open(filepath, 'w', encoding='utf-8') as f:
                f.write(doc_content)
            
            logger.info(f"✅ Text document created: {filepath}")
            return str(filepath)
            
        except Exception as e:
            logger.error(f"❌ Failed to create text document: {e}")
            # Return a temporary file path even if creation fails
            return str(self.base_output_dir / "documents" / f"failed_creation_{execution_id}.txt")
    
    def _create_speaker_notes(self, content: str, task: str, execution_id: str) -> str:
        """Create speaker notes for presentations"""
        try:
            notes_content = f"""
SPEAKER NOTES
=============

Presentation: {self._extract_title_from_task(task)}
Generated: {datetime.now().strftime('%B %d, %Y')}

CONTENT OVERVIEW:
{content}

PRESENTATION TIPS:
- Speak clearly and maintain eye contact with the audience
- Use the content as a guide, but feel free to elaborate with your own expertise
- Allow time for questions at the end
- Estimated duration: 15-30 minutes depending on audience interaction

TECHNICAL NOTES:
- This presentation was generated by AgentFlow AI
- Content is based on the latest available information
- Verify any specific data points before presenting to ensure accuracy
"""
            
            filename = f"Speaker_Notes_{execution_id}.txt"
            filepath = self.base_output_dir / "presentations" / filename
            
            with open(filepath, 'w', encoding='utf-8') as f:
                f.write(notes_content)
            
            return str(filepath)
            
        except Exception as e:
            logger.error(f"❌ Failed to create speaker notes: {e}")
            return ""
    
    def _create_metadata_file(self, content: str, task: str, execution_id: str, deliverables: List[str]) -> str:
        """Create metadata file with execution details"""
        try:
            metadata = {
                "execution_id": execution_id,
                "task": task,
                "generated_at": datetime.now().isoformat(),
                "content_length": len(content),
                "deliverables": deliverables,
                "generator": "AgentFlow AI System",
                "version": "1.0"
            }
            
            filename = f"metadata_{execution_id}.json"
            filepath = self.base_output_dir / "data" / filename
            
            with open(filepath, 'w', encoding='utf-8') as f:
                json.dump(metadata, f, indent=2)
            
            return str(filepath)
            
        except Exception as e:
            logger.error(f"❌ Failed to create metadata file: {e}")
            return ""
    
    def _structure_research_report(self, content: str, task: str) -> str:
        """Structure content as a research report"""
        title = self._extract_title_from_task(task)
        
        structured_content = f"""
{title}
{'=' * len(title)}

EXECUTIVE SUMMARY
-----------------
{content[:500]}...

DETAILED ANALYSIS
-----------------
{content}

CONCLUSIONS
-----------
Based on the analysis above, key findings include comprehensive insights into the subject matter.

RECOMMENDATIONS
---------------
Further research and implementation strategies should be considered based on these findings.

---
Report generated by AgentFlow AI
Date: {datetime.now().strftime('%B %d, %Y')}
"""
        return structured_content
    
    def _structure_data_analysis(self, content: str, task: str) -> str:
        """Structure content as a data analysis report"""
        title = self._extract_title_from_task(task)
        
        structured_content = f"""
{title}
{'=' * len(title)}

ANALYSIS OVERVIEW
-----------------
{content[:300]}...

KEY FINDINGS
------------
{content}

DATA INSIGHTS
-------------
The analysis reveals important patterns and trends in the data.

VISUALIZATIONS
--------------
[Charts and graphs would be displayed here in a full implementation]

CONCLUSIONS
-----------
Based on the data analysis, we can draw several important conclusions.

---
Analysis generated by AgentFlow AI
Date: {datetime.now().strftime('%B %d, %Y')}
"""
        return structured_content
    
    def _create_html_analysis_report(self, content: str, task: str, execution_id: str) -> str:
        """Create HTML analysis report"""
        try:
            title = self._extract_title_from_task(task)
            
            html_content = f"""
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <style>
        body {{ font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif; margin: 0; padding: 20px; background: #f8f9fa; }}
        .container {{ max-width: 1200px; margin: 0 auto; background: white; padding: 40px; border-radius: 10px; box-shadow: 0 4px 20px rgba(0,0,0,0.1); }}
        h1 {{ color: #2c3e50; border-bottom: 4px solid #3498db; padding-bottom: 15px; }}
        h2 {{ color: #34495e; margin-top: 30px; }}
        p {{ line-height: 1.8; color: #555; }}
        .highlight {{ background: #e8f4fd; padding: 15px; border-left: 4px solid #3498db; margin: 20px 0; }}
        .footer {{ margin-top: 40px; padding-top: 20px; border-top: 1px solid #eee; color: #888; text-align: center; }}
    </style>
</head>
<body>
    <div class="container">
        <h1>{title}</h1>
        <div class="highlight">
            <strong>Generated by AgentFlow AI</strong><br>
            Date: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}<br>
            Execution ID: {execution_id}
        </div>
        
        <div>
            {content.replace(chr(10), '<br>').replace('---', '<hr>')}
        </div>
        
        <div class="footer">
            This report was automatically generated by AgentFlow's AI system.
        </div>
    </div>
</body>
</html>
"""
            
            filename = f"Analysis_Report_{execution_id}.html"
            filepath = self.base_output_dir / "reports" / filename
            
            with open(filepath, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            logger.info(f"✅ HTML analysis report created: {filepath}")
            return str(filepath)
            
        except Exception as e:
            logger.error(f"❌ Failed to create HTML analysis report: {e}")
            return self._create_text_document(content, task, execution_id)
    
    def _create_text_report(self, content: str, task: str, execution_id: str) -> str:
        """Create text-based report"""
        filename = f"Report_{execution_id}.txt"
        filepath = self.base_output_dir / "reports" / filename
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(content)
        
        return str(filepath)
    
    def _create_pdf_report(self, content: str, task: str, execution_id: str) -> str:
        """Create PDF report"""
        if not FPDF:
            return self._create_text_report(content, task, execution_id)
        
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font('Arial', 'B', 16)
        
        title = self._extract_title_from_task(task)
        pdf.cell(0, 10, title, ln=True, align='C')
        pdf.ln(10)
        
        pdf.set_font('Arial', '', 11)
        
        # Add content
        lines = content.split('\n')
        for line in lines:
            if line.strip():
                pdf.cell(0, 6, line.encode('latin-1', 'replace').decode('latin-1'), ln=True)
        
        filename = f"Report_{execution_id}.pdf"
        filepath = self.base_output_dir / "reports" / filename
        pdf.output(str(filepath))
        
        return str(filepath)
    
    def get_deliverables_info(self) -> Dict[str, Any]:
        """Get information about created deliverables"""
        info = {
            "base_directory": str(self.base_output_dir),
            "subdirectories": {
                "presentations": str(self.base_output_dir / "presentations"),
                "documents": str(self.base_output_dir / "documents"),
                "reports": str(self.base_output_dir / "reports"),
                "data": str(self.base_output_dir / "data")
            },
            "total_files": 0,
            "file_types": {}
        }
        
        try:
            for subdir in info["subdirectories"].values():
                if os.path.exists(subdir):
                    files = os.listdir(subdir)
                    info["total_files"] += len(files)
                    
                    for file in files:
                        ext = os.path.splitext(file)[1]
                        info["file_types"][ext] = info["file_types"].get(ext, 0) + 1
        except Exception as e:
            logger.error(f"❌ Failed to get deliverables info: {e}")
        
        return info


    def create_basic_deliverables(self, task: str, execution_id: str) -> List[str]:
        """Create basic deliverables without AI when no API client is available"""
        created_files = []
        
        try:
            logger.info(f"📁 Creating basic deliverables for task: {task[:100]}...")
            
            # Create basic text document
            doc_content = f"""Basic Deliverable Document
Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
Execution ID: {execution_id}

Task: {task}

This is a basic deliverable created without AI assistance.
The system is configured to create actual files even when AI services are unavailable.

Content Summary:
- Task received and processed
- Basic deliverable structure created
- File generation completed successfully

Next Steps:
1. Review the generated content
2. Customize as needed for your specific requirements
3. Use this as a starting point for further development

Generated by AgentFlow - AI Workflow Automation Platform
"""
            
            # Save as text file
            doc_file = self.base_output_dir / "documents" / f"Basic_Document_{execution_id}.txt"
            doc_file.parent.mkdir(parents=True, exist_ok=True)
            doc_file.write_text(doc_content, encoding='utf-8')
            created_files.append(str(doc_file))
            logger.info(f"✅ Created basic document: {doc_file}")
            
            # Create metadata file
            metadata = {
                "execution_id": execution_id,
                "task": task,
                "created_at": datetime.now().isoformat(),
                "type": "basic_deliverable",
                "ai_used": False,
                "files_created": 1,
                "status": "completed_without_ai"
            }
            
            metadata_file = self.base_output_dir / "data" / f"metadata_{execution_id}.json"
            metadata_file.parent.mkdir(parents=True, exist_ok=True)
            metadata_file.write_text(json.dumps(metadata, indent=2), encoding='utf-8')
            created_files.append(str(metadata_file))
            logger.info(f"✅ Created metadata file: {metadata_file}")
            
            logger.info(f"🎉 Created {len(created_files)} basic deliverable files")
            return created_files
            
        except Exception as e:
            logger.error(f"❌ Failed to create basic deliverables: {e}")
            return []

